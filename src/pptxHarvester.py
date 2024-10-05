import os
from os import listdir
from os.path import isfile, join
import zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
import re
import nltk
from nltk.stem.porter import *
from nltk.stem.snowball import SnowballStemmer
porterStemmer = PorterStemmer()
from PIL import Image, ImageOps
import io
######################################################################################
paths_ = json.load(open("%s%s" % ("resources\\","paths.json")))
######################################################################################
def get_shape_location(slide, shape):
	# Check if the shape is present on the slide
	if shape not in slide.shapes:
		return None
	# Get the position (left and top) of the shape
	left = shape.left
	top = shape.top
	return left, top
######################################################################################
def wordStemmer(split_rawTxt,stopwords):
	acceptedChars = [
		"a","b","c","d","e","f","g","h",
		"i","j","k","l","m","n","o","p",
		"q","r","s","t","u","v","w","x","y","z"
		]
	words2stem = []
	for txt in split_rawTxt:
		len_txt = int(len(list(txt)))
		if len_txt >=3:
			if txt.lower() not in stopwords:
				alphaCharCount = 0
				for t in list(txt):
					if t.lower() in acceptedChars:
						alphaCharCount+=1
				if float(alphaCharCount+1/len_txt+1) >= .75:
					words2stem.append(txt.lower())
	#tokens = nltk.word_tokenize(' '.join(words2stem))
	stems_ = [porterStemmer.stem(word) for word in words2stem]
	stems_output = []
	terms_output = []
	for stem,term in zip(stems_,words2stem):
		if len(list(stem)) >=3:
			if stem not in stopwords:
				stems_output.append(stem)
				terms_output.append(term)
	return stems_output,terms_output;
######################################################################################
def extract_content(pptx_file,getTxt,getIm,fId):
	text_content = []
	image_content = []
	presentation = Presentation(pptx_file)
	slideId = 0
	for slide in presentation.slides:
		imCount=0
		txtCount = 0
		for shape in slide.shapes:
			parse_shape = str(shape).split(".")
			if getTxt == True:
				if shape.has_text_frame:
					txtCount+=1
					shapeLocation = get_shape_location(slide, shape)
					print(shapeLocation)
					text_frame = shape.text_frame
					for paragraph in text_frame.paragraphs:
						for run in paragraph.runs:
							text = run.text.strip()
							font_size = run.font.size.pt if run.font.size is not None else "Unknown"
							font_name = run.font.name if run.font.name is not None else "Unknown"
							
							text_with_font_info = f"text: {text}, font_size: {font_size}, font_name: {font_name}"
							text_content.append({
								"slideId":slideId,"txtId":txtCount,
								"text": text, "font_size": font_size, "font_name": font_name,
								#"font_uppercase":text.isupper(),
								"font_bold": run.font.bold, "font_italic": run.font.italic, "font_underline": run.font.underline,
								"shape_left":shapeLocation[0],"shape_top":shapeLocation[1]
								})
			if getIm == True:
				#if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
				if shape.shape_type == 13:
					imCount += 1
					imageShape = shape.image
					shapeLocation = get_shape_location(slide, shape)
					imId = str(fId)+"_"+str(slideId)+"_"+str(imCount)

					imWidth = shape.width
					imHeight = shape.height
					max_dimension = 500
					if imWidth > imHeight:
						new_width = max_dimension
						new_height = int(imHeight * (max_dimension / imWidth))
					else:
						new_height = max_dimension
						new_width = int(imWidth * (max_dimension / imHeight))
					try:
						pil_image = Image.open(io.BytesIO(imageShape.blob))

						imageResized = pil_image.resize((new_width,new_height))
						imageResized = imageResized.convert("RGB")
						imageResized.save(os.path.join(paths_["content"]["images"], f"im_{imId}.jpg"))
						image_content.append(
							{
								"slideId":slideId,"imageId":imId,
								"shape_left":shapeLocation[0],"shape_top":shapeLocation[1]
							}
						)
					except OSError as e:
						print(f"Error processing {imId}: {e}")
						pass
		slideId+=1
	return text_content,image_content;
######################################################################################
######################################################################################
textMiningResources_ = json.load(open("%s%s" % (paths_["00_resources"],"TextMining.json")))
stopwords_ = textMiningResources_["nltk_resources"]["stopwords"]
ledger_ = json.load(open("%s%s" % (r"resources\\","ledger.json")))
######################################################################################
if __name__ == "__main__":
	files_ = [f for f in listdir(paths_["projectDir"]) if isfile(join(paths_["projectDir"], f))]
	for f in files_:
		parse_f  = f.split(".")
		if parse_f[-1] == "pptx":
			id = None
			for key,item in ledger_.items():
				if parse_f[0] == item["name"] and parse_f[-1] == item["ext"]:
					id = key
					break
			try:
				input_file =  "%s%s" % (paths_["projectDir"],f)
				extracted_content = extract_content(input_file,True,True,id)
				manifest_ = {}
				for content_ in extracted_content[0]:
					slideId = content_["slideId"]
					parse_txt = re.split(r"[   ‘'’“”  :€™;,<>/\|*_ ~*&^%#@=()]", content_["text"])
					stemmer_ = wordStemmer(parse_txt,stopwords_)
					#print(stemmer_[1])
					if slideId not in list(manifest_.keys()):
						manifest_[slideId] = {
							"text":{content_["txtId"]:{
								"text":content_["text"],
								"text_no_stopwords":stemmer_[1],
								"stems_no_stopwords":stemmer_[0],
								"font":content_["font_name"],
								"size":content_["font_size"],
								"font_bold": content_["font_bold"],
								"font_italic": content_["font_italic"],
								"font_underline": content_["font_underline"],
								"shape_left":content_["shape_left"],
								"shape_top":content_["shape_top"]}
							},
							"image":{}
						}
					else:
						manifest_[slideId]["text"][content_["txtId"]] = {
							"text":content_["text"],
							"text_no_stopwords":stemmer_[1],
							"stems_no_stopwords":stemmer_[0],
							"font":content_["font_name"],
							"size":content_["font_size"],
							"font_bold": content_["font_bold"],
							"font_italic": content_["font_italic"],
							"font_underline": content_["font_underline"],
							"shape_left":content_["shape_left"],
							"shape_top":content_["shape_top"]
						}
						
				for content_ in extracted_content[1]:
					slideId = content_["slideId"]
					if slideId not in list(manifest_.keys()):
						manifest_[slideId] = {
							"image":{content_["imageId"]:{
								"shape_left":content_["shape_left"],
								"shape_top":content_["shape_top"]
								}
							}
						}
					else:
						manifest_[slideId]["image"][content_["imageId"]] = {
							"shape_left":content_["shape_left"],
							"shape_top":content_["shape_top"]
						}
			except Exception as e:
				print(e)
				pass
######################################################################################
######################################################################################
			#ensure_ascii=False
			with open(str(
				"%s%s%s" % (paths_["content"]["manifests"],parse_f[0],"_manifest.json")
				), "w", encoding='utf-8') as json_manifest:
				json_manifest.write(json.dumps(manifest_, indent=4))